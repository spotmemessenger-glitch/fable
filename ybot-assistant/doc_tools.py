"""Ybot's document skill — mirrors Anthropic's document-skills (PDF/DOCX/PPTX/XLSX)
as a real Python capability instead of a Claude-only skill file.

The brain (via ai_client, free Gemini route) plans STRUCTURED CONTENT as JSON; this
module deterministically builds the actual file with the right library — the model
never touches binary formats or raw XML, so output is always a valid, openable file.
"""
from __future__ import annotations

import json
import re
import time
from pathlib import Path

import ai_client

_ROOT = Path(__file__).resolve().parent
_JOBS = _ROOT / "creations" / "documents"


def _slug(text: str) -> str:
    s = re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-")
    return s[:40] or "document"


def _extract_json(text: str) -> dict:
    text = re.sub(r"^```(?:json)?\n?", "", text.strip())
    text = re.sub(r"\n?```$", "", text).strip()
    m = re.search(r"\{.*\}", text, re.DOTALL)
    return json.loads(m.group(0) if m else text)


_SYSTEM = (
    "You plan document CONTENT ONLY — never markup, never code. Given a task and a "
    "target format, respond with ONLY this JSON (no markdown fences, no prose):\n"
    '{"format": "docx|pptx|xlsx|pdf", "title": "...", "content": <shape below>}\n\n'
    "Shape of `content` per format:\n"
    '- docx/pdf: {"paragraphs": [{"style": "Title|Heading 1|Heading 2|Normal", "text": "..."}]}\n'
    '- pptx: {"slides": [{"title": "...", "bullets": ["...", "..."]}]}\n'
    '- xlsx: {"sheets": [{"name": "...", "rows": [["col1","col2"], ["v1","v2"]]}]} '
    "(first row is the header)\n\n"
    "Write real, complete, useful content for the task — no placeholders, no "
    '"[insert X here]". Pick whichever format best fits the task if the user did not '
    "specify one."
)


def _plan(task: str) -> dict:
    msg = ai_client.client().messages.create(
        model=ai_client.TEXT_MODEL, max_tokens=4096, system=_SYSTEM,
        messages=[{"role": "user", "content": f"Task: {task}\n\nPlan the document now."}],
    )
    text = "".join(b.text for b in msg.content if getattr(b, "type", "") == "text").strip()
    return _extract_json(text)


def _build_docx(path: Path, title: str, content: dict) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=0)
    for p in content.get("paragraphs", []):
        style = p.get("style", "Normal")
        if style.startswith("Heading") or style == "Title":
            level = 1 if style == "Title" else int(style[-1]) if style[-1].isdigit() else 1
            doc.add_heading(p.get("text", ""), level=level)
        else:
            doc.add_paragraph(p.get("text", ""))
    doc.save(str(path))


def _build_pptx(path: Path, title: str, content: dict) -> None:
    from pptx import Presentation

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    for s in content.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = s.get("title", "")
        body = slide.placeholders[1].text_frame
        bullets = s.get("bullets", [])
        if bullets:
            body.text = bullets[0]
            for b in bullets[1:]:
                para = body.add_paragraph()
                para.text = b
    prs.save(str(path))


def _build_xlsx(path: Path, title: str, content: dict) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)
    for sheet in content.get("sheets", [{"name": "Sheet1", "rows": []}]):
        ws = wb.create_sheet(sheet.get("name", "Sheet1")[:31])
        for row in sheet.get("rows", []):
            ws.append(row)
    if not wb.sheetnames:
        wb.create_sheet("Sheet1")
    wb.save(str(path))


def _build_pdf(path: Path, title: str, content: dict) -> None:
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    styles = getSampleStyleSheet()
    story = [Paragraph(title, styles["Title"]), Spacer(1, 12)]
    for p in content.get("paragraphs", []):
        style_name = p.get("style", "Normal")
        style = styles.get(style_name if style_name in styles else "Normal", styles["Normal"])
        if style_name.startswith("Heading"):
            style = styles.get("Heading2", styles["Heading1"])
        story.append(Paragraph(p.get("text", ""), style))
        story.append(Spacer(1, 8))
    SimpleDocTemplate(str(path)).build(story)


_BUILDERS = {"docx": _build_docx, "pptx": _build_pptx, "xlsx": _build_xlsx, "pdf": _build_pdf}


def create_document(task: str, progress=None) -> dict:
    """Plain-English task -> a real, openable PDF/DOCX/PPTX/XLSX file."""
    def note(msg: str) -> None:
        if progress:
            try:
                progress(msg)
            except Exception:
                pass

    note("Planning the document content...")
    plan = _plan(task)
    fmt = (plan.get("format") or "docx").lower()
    if fmt not in _BUILDERS:
        return {"ok": False, "errors": [f"unknown format {fmt!r} from the brain"]}
    title = plan.get("title") or task[:60]

    _JOBS.mkdir(parents=True, exist_ok=True)
    job = _JOBS / f"{_slug(title)}-{time.strftime('%Y%m%d-%H%M%S')}"
    job.mkdir(parents=True, exist_ok=True)
    out_path = job / f"{_slug(title)}.{fmt}"

    note(f"Writing the .{fmt} file...")
    try:
        _BUILDERS[fmt](out_path, title, plan.get("content", {}))
    except Exception as e:
        return {"ok": False, "job_dir": str(job), "errors": [f"{type(e).__name__}: {e}"]}

    return {
        "ok": out_path.exists(),
        "job_dir": str(job),
        "file": str(out_path),
        "format": fmt,
        "title": title,
        "errors": [] if out_path.exists() else ["file was not created"],
    }


if __name__ == "__main__":
    import sys

    task = " ".join(sys.argv[1:]) or "a one-page PDF summary of the water cycle"
    print(json.dumps(create_document(task, progress=print), indent=2))
